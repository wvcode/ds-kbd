from pptx import Presentation
from copy import deepcopy
import os


def copy_slide(prs, slide):
    SLD_LAYOUT = 0
    slide_layout = prs.slide_layouts[SLD_LAYOUT]

    new_slide = prs.slides.add_slide(slide_layout)

    imgDict = {}

    for shp in slide.shapes:
        print(shp.name)
        if "Google" in shp.name:
            # save image
            with open(shp.name + ".jpg", "wb") as f:
                try:
                    f.write(shp.image.blob)
                    imgDict[shp.name + ".jpg"] = [
                        shp.left,
                        shp.top,
                        shp.width,
                        shp.height,
                    ]
                except:
                    el = shp.element
                    newel = deepcopy(el)
                    new_slide.shapes._spTree.insert_element_before(newel, "p:extLst")

        else:
            el = shp.element
            newel = deepcopy(el)
            new_slide.shapes._spTree.insert_element_before(newel, "p:extLst")

    for k, v in imgDict.items():
        new_slide.shapes.add_picture(k, v[0], v[1], v[2], v[3])
        os.remove(k)

    return new_slide


def create_slide(prs, slide_to_copy, item):
    slide = copy_slide(prs, slide_to_copy)
    for shape in slide.shapes:
        if shape.has_text_frame:
            if shape.text == "id":
                shape.text = item["label"]
            elif shape.text == "nivel":
                shape.text = item["nivel"]
            elif shape.text == "tipo":
                shape.text = item["tipo"]
            elif shape.text == "area":
                shape.text = item["area"]
            elif shape.text == "questao":
                shape.text = item["resultado"]
    return slide


def create_presentation(slides, template_filename, presentation_filename):
    # Carregue a apresentação existente
    prs = Presentation(template_filename)
    slide_to_copy = prs.slides[0]

    for idx, slide in slides.iterrows():
        create_slide(prs, slide_to_copy, slide)

    # Salvar a apresentação em um arquivo
    prs.save(presentation_filename)
